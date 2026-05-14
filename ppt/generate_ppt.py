"""
Generate presentation PPT for sushepro (Dormitory Management System).
Minimalist style, 16:9 widescreen, MD3 color scheme.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import os

# ── Design Tokens (from head.jsp MD3) ──────────────────────────────
PRIMARY = RGBColor(0x00, 0x58, 0xBC)
PRIMARY_DARK = RGBColor(0x00, 0x44, 0x93)
ON_SURFACE = RGBColor(0x19, 0x1C, 0x1D)
ON_SURFACE_VARIANT = RGBColor(0x41, 0x47, 0x55)
SURFACE = RGBColor(0xF8, 0xF9, 0xFA)
SURFACE_CONTAINER = RGBColor(0xED, 0xEE, 0xEF)
SURFACE_LOWEST = RGBColor(0xFF, 0xFF, 0xFF)
OUTLINE_VARIANT = RGBColor(0xC1, 0xC6, 0xD7)
OUTLINE = RGBColor(0x71, 0x77, 0x86)
SECONDARY = RGBColor(0x00, 0x67, 0x7D)
ERROR = RGBColor(0xBA, 0x1A, 0x1A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = 'Microsoft YaHei'
FONT_EN = 'Inter'

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helper Functions ────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    """Add a rectangle shape, optionally with fill and border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape


def add_rounded_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        shape.line.width = Pt(1)
    # Adjust corner radius
    shape.adjustments[0] = 0.1
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=Pt(14),
                 color=ON_SURFACE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name=FONT, anchor=MSO_ANCHOR.TOP, line_spacing=1.3):
    """Add a text box with single text run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing:
        p.line_spacing = Pt(font_size.pt * line_spacing)
    tf.paragraphs[0].font.east_asian_font = font_name
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=Pt(13),
                      color=ON_SURFACE, line_spacing=1.5, anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple paragraphs. lines is list of (text, bold, font_size_override, color_override)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, fs, clr = line_data, False, font_size, color
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            fs = line_data[2] if len(line_data) > 2 else font_size
            clr = line_data[3] if len(line_data) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = fs
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = FONT
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(2)
        p.space_before = Pt(2)
        p.line_spacing = Pt(fs.pt * line_spacing)
        p.font.east_asian_font = FONT
    return txBox


def add_badge(slide, left, top, text, color=PRIMARY):
    """Add a pill-shaped badge with text."""
    # Estimate width from text length
    est_w = Inches(0.4 + len(text) * 0.14)
    h = Inches(0.32)
    shape = add_rounded_rect(slide, left, top, est_w, h, fill_color=color)
    shape.adjustments[0] = 0.5  # fully rounded pill
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    p.font.east_asian_font = FONT
    return shape


def add_accent_header(slide, left, top, width, title_text, subtitle_text=None):
    """Add a section header with blue left accent bar."""
    # Accent bar
    add_rect(slide, left, top, Inches(0.06), Inches(0.55), fill_color=PRIMARY)
    # Title
    add_text_box(slide, left + Inches(0.18), top - Inches(0.04), width - Inches(0.18), Inches(0.45),
                 title_text, font_size=Pt(22), bold=True, color=ON_SURFACE)
    if subtitle_text:
        add_text_box(slide, left + Inches(0.18), top + Inches(0.4), width - Inches(0.18), Inches(0.3),
                     subtitle_text, font_size=Pt(11), color=ON_SURFACE_VARIANT)


def add_card(slide, left, top, width, height, title, body_lines, accent_color=PRIMARY, title_size=Pt(15)):
    """Add a content card with left accent border."""
    # Card background
    card = add_rounded_rect(slide, left, top, width, height, fill_color=SURFACE_LOWEST, border_color=OUTLINE_VARIANT)
    # Left accent bar
    add_rect(slide, left + Inches(0.02), top + Inches(0.1), Inches(0.05), height - Inches(0.2),
             fill_color=accent_color)
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.08), width - Inches(0.35), Inches(0.35),
                 title, font_size=title_size, bold=True, color=ON_SURFACE)
    # Body
    add_multiline_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.35),
                      height - Inches(0.55), body_lines, font_size=Pt(11), color=ON_SURFACE_VARIANT)


def add_page_number(slide, num):
    add_text_box(slide, SLIDE_W - Inches(0.7), SLIDE_H - Inches(0.45), Inches(0.5), Inches(0.3),
                 str(num), font_size=Pt(10), color=OUTLINE, alignment=PP_ALIGN.RIGHT)


def add_connector(slide, start_left, start_top, end_left, end_top, color=OUTLINE_VARIANT):
    """Add a straight connector line between two points."""
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, start_left, start_top, end_left, end_top)
    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)
    return connector


def add_arrow_right(slide, left, top, width, color=PRIMARY):
    """Add a right-pointing arrow."""
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, Inches(0.22))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = color
    arrow.line.fill.background()
    return arrow


# ══════════════════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════════════════════════
def slide_cover():
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, SURFACE)

    # Top accent bar
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    # Center content
    cx, cy = SLIDE_W / 2, SLIDE_H / 2

    # Icon placeholder (rounded square with building symbol)
    icon = add_rounded_rect(slide, cx - Inches(0.45), cy - Inches(1.9), Inches(0.9), Inches(0.9),
                            fill_color=PRIMARY)
    tf = icon.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = "⌂"  # ⌂ house/building
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Main title
    add_text_box(slide, cx - Inches(3), cy - Inches(0.8), Inches(6), Inches(0.7),
                 "宿舍管理系统", font_size=Pt(40), bold=True, color=ON_SURFACE,
                 alignment=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, cx - Inches(3), cy - Inches(0.15), Inches(6), Inches(0.5),
                 "基于 SSM 框架的智能宿舍管理与 AI 推荐平台", font_size=Pt(16),
                 color=ON_SURFACE_VARIANT, alignment=PP_ALIGN.CENTER)

    # Divider line
    add_rect(slide, cx - Inches(1.2), cy + Inches(0.45), Inches(2.4), Inches(0.02),
             fill_color=OUTLINE_VARIANT)

    # Tech badges
    badges = ["Spring 6", "Spring MVC", "MyBatis", "MySQL 8.0", "DeepSeek AI", "ECharts", "Tailwind CSS"]
    total_badges_w = sum(0.4 + len(b) * 0.14 for b in badges) + (len(badges) - 1) * 0.15
    start_x = cx - Inches(total_badges_w / 2)
    cur_x = start_x
    for b in badges:
        add_badge(slide, cur_x, cy + Inches(0.7), b)
        cur_x += Inches(0.4 + len(b) * 0.14 + 0.15)

    # Bottom info
    add_text_box(slide, Inches(0.8), SLIDE_H - Inches(0.55), Inches(4), Inches(0.3),
                 "2026 · 毕业设计项目", font_size=Pt(11), color=OUTLINE)

    add_page_number(slide, 1)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 2: PROJECT OVERVIEW
# ══════════════════════════════════════════════════════════════════════
def slide_overview():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "项目概述", "Project Overview")

    # Two-column layout: Background / Goals
    col_w = Inches(5.7)
    col_h = Inches(2.4)

    add_card(slide, Inches(0.7), Inches(1.3), col_w, col_h,
             "项目背景",
             [("高校宿舍管理面临学生信息分散、调换审批繁琐、报修跟踪困难等问题。",
               False, Pt(12)),
              ("传统人工管理方式效率低、易出错，亟需数字化系统提升管理效能。",
               False, Pt(12)),
              ("结合 AI 技术实现智能宿舍推荐与客服，为管理者提供数据驱动的决策支持。",
               False, Pt(12))])

    add_card(slide, Inches(6.9), Inches(1.3), col_w, col_h,
             "项目目标",
             [("构建一套完整的宿舍管理数字化平台，覆盖学生入住、调换、报修全流程。",
               False, Pt(12)),
              ("实现双角色权限体系：管理员进行全局管理，学生自助办理宿舍相关业务。",
               False, Pt(12)),
              ("引入 DeepSeek AI 大模型，提供智能宿舍匹配推荐与 24×7 在线客服。",
               False, Pt(12))],
             accent_color=SECONDARY)

    # Features row
    add_text_box(slide, Inches(0.7), Inches(4.1), Inches(4), Inches(0.35),
                 "系统特性", font_size=Pt(14), bold=True, color=ON_SURFACE)

    features = [
        ("", "双角色权限", "Admin / Student\n精细化权限控制"),
        ("", "生活习惯问卷", "六维生活习惯采集\n科学匹配依据"),
        ("", "AI 智能推荐", "DeepSeek 大模型驱动\n最优宿舍匹配"),
        ("", "数据可视化", "ECharts 雷达图\n六维匹配分析"),
    ]
    card_w = Inches(2.85)
    start_x = Inches(0.7)
    gap = Inches(0.18)
    for i, (icon, title, desc) in enumerate(features):
        x = start_x + i * (card_w + gap)
        y = Inches(4.5)
        card = add_rounded_rect(slide, x, y, card_w, Inches(1.6), fill_color=SURFACE_LOWEST,
                                border_color=OUTLINE_VARIANT)
        # Icon circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.15), y + Inches(0.2),
                                         Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(0.4), Inches(0.4),
                     icon, font_size=Pt(14), color=WHITE, alignment=PP_ALIGN.CENTER,
                     anchor=MSO_ANCHOR.MIDDLE, font_name=FONT_EN)
        add_text_box(slide, x + Inches(0.68), y + Inches(0.18), card_w - Inches(0.85), Inches(0.35),
                     title, font_size=Pt(13), bold=True, color=ON_SURFACE)
        add_text_box(slide, x + Inches(0.68), y + Inches(0.7), card_w - Inches(0.85), Inches(0.8),
                     desc, font_size=Pt(10), color=ON_SURFACE_VARIANT)

    add_page_number(slide, 2)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 3: TECH ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════
def slide_architecture():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "技术架构", "Architecture")

    # Architecture layers - centered stacked boxes
    layers = [
        ("表现层", "JSP + Tailwind CSS + ECharts\nMaterial Design 3 设计语言", Inches(2.8)),
        ("控制层", "Spring MVC DispatcherServlet\nRESTful 风格路由，拦截器校验权限", Inches(2.8)),
        ("业务层", "Spring IoC + AOP 事务管理\n8 个 Service 接口，模块化设计", Inches(2.8)),
        ("持久层", "MyBatis + Druid 连接池\n7 个 Mapper 接口 + XML 映射文件", Inches(2.8)),
        ("数据层", "MySQL 8.0 关系型数据库\n7 张业务表，完整外键约束", Inches(2.8)),
    ]

    layer_h = Inches(0.72)
    layer_w = Inches(4.5)
    layer_gap = Inches(0.08)
    cx = SLIDE_W / 2

    y_start = Inches(1.25)
    for i, (name, desc, w) in enumerate(layers):
        y = y_start + i * (layer_h + layer_gap)
        lx = cx - layer_w / 2
        box = add_rounded_rect(slide, lx, y, layer_w, layer_h, fill_color=SURFACE_LOWEST,
                               border_color=PRIMARY if i == 0 else OUTLINE_VARIANT)
        # Layer name (left part)
        add_text_box(slide, lx + Inches(0.2), y + Inches(0.08), Inches(1.2), Inches(0.3),
                     name, font_size=Pt(12), bold=True, color=PRIMARY if i == 0 else ON_SURFACE)
        add_text_box(slide, lx + Inches(0.2), y + Inches(0.38), Inches(1.2), Inches(0.25),
                     "", font_size=Pt(9), color=ON_SURFACE_VARIANT)
        # Description
        add_text_box(slide, lx + Inches(1.5), y + Inches(0.12), layer_w - Inches(1.7), layer_h - Inches(0.2),
                     desc, font_size=Pt(10), color=ON_SURFACE_VARIANT)
        # Connector arrow between layers
        if i < len(layers) - 1:
            arrow_y = y + layer_h
            add_connector(slide, cx, arrow_y, cx, arrow_y + layer_gap, OUTLINE_VARIANT)

    # AI side chain
    ai_y = y_start + 1 * (layer_h + layer_gap)
    ai_x = cx - layer_w / 2 - Inches(2.3)
    ai_box = add_rounded_rect(slide, ai_x, ai_y, Inches(1.9), Inches(1.5), fill_color=SURFACE_LOWEST,
                              border_color=SECONDARY)
    add_text_box(slide, ai_x + Inches(0.15), ai_y + Inches(0.1), Inches(1.6), Inches(0.3),
                 "DeepSeek AI", font_size=Pt(12), bold=True, color=SECONDARY)
    add_multiline_box(slide, ai_x + Inches(0.15), ai_y + Inches(0.45), Inches(1.6), Inches(0.9),
                      [("智能宿舍推荐", False, Pt(10), ON_SURFACE_VARIANT),
                       ("AI 客服聊天", False, Pt(10), ON_SURFACE_VARIANT),
                       ("Apache HttpClient 5", False, Pt(10), ON_SURFACE_VARIANT)],
                      font_size=Pt(10))

    # Connect AI to Service layer
    add_connector(slide, ai_x + Inches(1.9), ai_y + Inches(0.75),
                  cx - layer_w / 2, ai_y + Inches(0.75), SECONDARY)

    # Right side tech detail columns
    right_x = cx + layer_w / 2 + Inches(0.5)
    col_w2 = Inches(2.6)
    col_data = [
        ("前端", ["JSP + JSTL 3.0", "Tailwind CSS CDN", "Material Symbols 图标", "ECharts 5.5 图表", "", ""]),
        ("后端", ["Spring 6.1.3", "Spring MVC", "MyBatis 3.5.16", "Druid 连接池", "Commons FileUpload", "", ""]),
        ("数据 & AI", ["MySQL 8.0", "DeepSeek Chat API", "FastJSON 2.0", "Log4j 日志", "Maven 构建", "Tomcat 部署"]),
    ]

    for i, (col_title, items) in enumerate(col_data):
        cx2 = right_x + i * (col_w2 + Inches(0.15))
        add_card(slide, cx2, y_start, col_w2, y_start + Inches(5.2) - y_start,
                 col_title, [(item, False, Pt(10)) for item in items if item],
                 accent_color=PRIMARY if i == 0 else (SECONDARY if i == 2 else ON_SURFACE),
                 title_size=Pt(13))

    add_page_number(slide, 3)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 4: MODULE OVERVIEW
# ══════════════════════════════════════════════════════════════════════
def slide_modules():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "系统功能模块", "Modules")

    # Admin card
    admin_items = [
        ("", "仪表盘", "数据概览与最近动态"),
        ("", "学生管理", "学生信息查看"),
        ("", "宿舍管理", "宿舍楼栋增删改查"),
        ("", "习惯问卷", "查看学生生活习惯"),
        ("", "分配管理", "手动分配宿舍"),
        ("", "调换审批", "审核调换申请"),
        ("", "报修处理", "处理维修工单"),
        ("", "AI 推荐", "智能宿舍匹配"),
    ]

    student_items = [
        ("", "个人仪表盘", "住宿状态一览"),
        ("", "个人信息", "查看学籍信息"),
        ("", "习惯问卷", "填写生活习惯"),
        ("", "宿舍分配", "查看分配结果"),
        ("", "调换申请", "提交调换请求"),
        ("", "报修申请", "提交维修工单"),
        ("", "AI 客服", "24×7 智能问答"),
        ("", "个人设置", "密码修改"),
    ]

    card_w2 = Inches(5.8)
    card_h2 = Inches(5.1)
    y2 = Inches(1.15)

    # Admin module card
    add_card(slide, Inches(0.7), y2, card_w2, card_h2, "管理员端功能", [],
             accent_color=PRIMARY, title_size=Pt(16))
    for j, (icon, title, desc) in enumerate(admin_items):
        row = j // 2
        col = j % 2
        ix = Inches(0.7) + Inches(0.25) + col * Inches(2.75)
        iy = y2 + Inches(0.55) + row * Inches(0.65)
        if j < 6:
            pass
        add_text_box(slide, ix, iy, Inches(0.35), Inches(0.35),
                     icon, font_size=Pt(14), color=PRIMARY, font_name=FONT_EN,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, ix + Inches(0.4), iy - Inches(0.02), Inches(1.9), Inches(0.22),
                     title, font_size=Pt(12), bold=True, color=ON_SURFACE)
        add_text_box(slide, ix + Inches(0.4), iy + Inches(0.22), Inches(1.9), Inches(0.22),
                     desc, font_size=Pt(9), color=ON_SURFACE_VARIANT)

    # Student module card
    add_card(slide, Inches(6.85), y2, card_w2, card_h2, "学生端功能", [],
             accent_color=SECONDARY, title_size=Pt(16))
    for j, (icon, title, desc) in enumerate(student_items):
        row = j // 2
        col = j % 2
        ix = Inches(6.85) + Inches(0.25) + col * Inches(2.75)
        iy = y2 + Inches(0.55) + row * Inches(0.65)
        add_text_box(slide, ix, iy, Inches(0.35), Inches(0.35),
                     icon, font_size=Pt(14), color=SECONDARY, font_name=FONT_EN,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, ix + Inches(0.4), iy - Inches(0.02), Inches(1.9), Inches(0.22),
                     title, font_size=Pt(12), bold=True, color=ON_SURFACE)
        add_text_box(slide, ix + Inches(0.4), iy + Inches(0.22), Inches(1.9), Inches(0.22),
                     desc, font_size=Pt(9), color=ON_SURFACE_VARIANT)

    # Bottom summary bar
    bar = add_rounded_rect(slide, Inches(0.7), SLIDE_H - Inches(1.0), Inches(11.9), Inches(0.5),
                           fill_color=PRIMARY)
    add_text_box(slide, Inches(0.7), SLIDE_H - Inches(1.0), Inches(11.9), Inches(0.5),
                 "系统共包含 7 张数据库表，支持 Admin / Student 双角色权限体系，覆盖宿舍管理全业务流程",
                 font_size=Pt(11), color=WHITE, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_page_number(slide, 4)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 5: ADMIN FEATURES
# ══════════════════════════════════════════════════════════════════════
def slide_admin():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "管理员核心功能", "Admin Features")

    # Three feature cards
    cards_data = [
        ("学生管理", "查看全部学生信息\n查看学生生活习惯问卷\n查看学生宿舍分配详情", ""),
        ("宿舍管理", "宿舍楼栋增删改查\n宿舍容量与入住管理\n按楼栋 / 类型筛选查看", ""),
        ("审批中心", "调换宿舍申请审批\n报修工单处理流转\n支持通过 / 拒绝 / 删除操作", ""),
    ]

    card_w3 = Inches(3.7)
    card_h3 = Inches(2.2)
    y3 = Inches(1.2)
    for i, (title, desc, icon) in enumerate(cards_data):
        x = Inches(0.7) + i * (card_w3 + Inches(0.25))
        add_card(slide, x, y3, card_w3, card_h3, title,
                 [(line, False, Pt(11)) for line in desc.split("\n")],
                 accent_color=PRIMARY, title_size=Pt(15))
        add_text_box(slide, x + card_w3 - Inches(0.7), y3 + Inches(0.1), Inches(0.5), Inches(0.5),
                     icon, font_size=Pt(28), color=OUTLINE_VARIANT, font_name=FONT_EN,
                     anchor=MSO_ANCHOR.MIDDLE)

    # Transfer workflow
    add_text_box(slide, Inches(0.7), Inches(3.7), Inches(4), Inches(0.35),
                 "宿舍调换审批流程", font_size=Pt(14), bold=True, color=ON_SURFACE)

    steps = ["学生提交\n调换申请", "管理员\n审核材料", "查看目标\n宿舍信息", "通过 / 拒绝\n填写意见", "通知学生\n完成调换"]
    step_w = Inches(2.2)
    y4 = Inches(4.2)

    for i, step_text in enumerate(steps):
        sx = Inches(0.7) + i * (step_w + Inches(0.18))
        box = add_rounded_rect(slide, sx, y4, step_w, Inches(1.0), fill_color=SURFACE_LOWEST,
                               border_color=PRIMARY)
        # Step number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, sx + step_w / 2 - Inches(0.2),
                                         y4 - Inches(0.15), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_EN
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, sx, y4 + Inches(0.3), step_w, Inches(0.6),
                     step_text, font_size=Pt(10), color=ON_SURFACE, alignment=PP_ALIGN.CENTER)

        if i < len(steps) - 1:
            add_arrow_right(slide, sx + step_w + Inches(0.02), y4 + Inches(0.4), Inches(0.16), PRIMARY)

    # Additional stats
    add_text_box(slide, Inches(0.7), Inches(5.5), Inches(11.9), Inches(0.3),
                 "仪表盘实时统计：学生总数 | 入住率 | 待分配人数 | 待处理报修 | 最近活动动态",
                 font_size=Pt(11), color=ON_SURFACE_VARIANT)

    add_page_number(slide, 5)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 6: STUDENT FEATURES
# ══════════════════════════════════════════════════════════════════════
def slide_student():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "学生核心功能", "Student Features")

    # 2x2 card grid
    s_cards = [
        ("生活习惯问卷", "",
         ["填写六维生活习惯信息", "睡眠时间 / 起床时间", "吸烟 / 饮酒习惯", "噪音容忍度 / 清洁习惯"],
         SECONDARY),
        ("宿舍分配查看", "",
         ["查看已分配的宿舍信息", "宿舍号 / 楼栋 / 房型", "入住状态实时更新", "支持查看同住室友"],
         PRIMARY),
        ("报修与调换", "",
         ["在线提交维修报修申请", "支持图片上传附件", "在线提交宿舍调换请求", "实时追踪处理进度"],
         SECONDARY),
        ("AI 智能客服", "",
         ["24×7 在线智能问答", "DeepSeek 大模型驱动", "宿舍相关业务咨询", "对话历史持久保存"],
         PRIMARY),
    ]

    card_w4 = Inches(5.7)
    card_h4 = Inches(2.55)
    start_y = Inches(1.15)
    gap_x4 = Inches(0.35)
    gap_y4 = Inches(0.25)

    for i, (title, icon, lines, accent) in enumerate(s_cards):
        row, col = i // 2, i % 2
        x = Inches(0.7) + col * (card_w4 + gap_x4)
        y = start_y + row * (card_h4 + gap_y4)

        add_card(slide, x, y, card_w4, card_h4, title, [], accent_color=accent, title_size=Pt(15))
        # Big icon in top-right
        add_text_box(slide, x + card_w4 - Inches(0.85), y + Inches(0.2), Inches(0.7), Inches(0.7),
                     icon, font_size=Pt(36), color=OUTLINE_VARIANT, font_name=FONT_EN,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_multiline_box(slide, x + Inches(0.2), y + Inches(0.55), card_w4 - Inches(0.35),
                          card_h4 - Inches(0.65),
                          [(line, False, Pt(11)) for line in lines],
                          font_size=Pt(11), color=ON_SURFACE_VARIANT)

    add_page_number(slide, 6)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 7: AI RECOMMENDATION
# ══════════════════════════════════════════════════════════════════════
def slide_ai_recommend():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11),
                       "AI 智能推荐系统", "DeepSeek-Powered Recommendation")

    # Left: workflow steps
    add_text_box(slide, Inches(0.7), Inches(1.15), Inches(3), Inches(0.35),
                 "推荐流程", font_size=Pt(15), bold=True, color=ON_SURFACE)

    steps_ai = [
        ("1", "选择学生", "管理员输入目标\n学生学号"),
        ("2", "采集数据", "系统自动收集学生\n习惯问卷 + 宿舍信息"),
        ("3", "构建 Prompt", "组装六维匹配参数\n发送至 DeepSeek API"),
        ("4", "解析结果", "提取 [SCORES] 标签\n解析六维评分"),
        ("5", "展示推荐", "推荐文本 + ECharts\n雷达图可视化"),
    ]

    step_h5 = Inches(0.95)
    for i, (num, title, desc) in enumerate(steps_ai):
        y = Inches(1.65) + i * (step_h5 + Inches(0.08))
        # Number circle
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), y + Inches(0.1),
                                         Inches(0.38), Inches(0.38))
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = FONT_EN
        p.alignment = PP_ALIGN.CENTER
        # Text
        add_text_box(slide, Inches(1.4), y + Inches(0.02), Inches(1.5), Inches(0.28),
                     title, font_size=Pt(13), bold=True, color=ON_SURFACE)
        add_text_box(slide, Inches(1.4), y + Inches(0.32), Inches(2.5), Inches(0.55),
                     desc, font_size=Pt(10), color=ON_SURFACE_VARIANT)
        # Connector
        if i < len(steps_ai) - 1:
            next_y = y + step_h5 + Inches(0.04)
            add_connector(slide, Inches(1.04), y + step_h5, Inches(1.04), next_y, OUTLINE_VARIANT)

    # Middle: 6 dimensions
    dim_x = Inches(4.5)
    dims = ["睡眠时间", "起床时间", "吸烟", "饮酒", "噪音容忍", "清洁习惯"]
    add_text_box(slide, dim_x, Inches(1.15), Inches(3), Inches(0.35),
                 "六维匹配维度", font_size=Pt(15), bold=True, color=ON_SURFACE)

    for i, dim in enumerate(dims):
        y = Inches(1.65) + i * Inches(0.52)
        bar_bg = add_rounded_rect(slide, dim_x, y, Inches(2.5), Inches(0.35),
                                  fill_color=SURFACE_CONTAINER)
        # Score bar (varying widths)
        scores = [0.85, 0.78, 0.92, 0.88, 0.75, 0.80]
        bar_fill = add_rounded_rect(slide, dim_x, y, Inches(2.5 * scores[i]), Inches(0.35),
                                    fill_color=PRIMARY)
        dim_name = add_text_box(slide, dim_x + Inches(0.12), y + Inches(0.02), Inches(2.3), Inches(0.3),
                                dim, font_size=Pt(11), bold=True, color=WHITE)
        score_label = add_text_box(slide, dim_x + Inches(1.2), y + Inches(0.02), Inches(1.2), Inches(0.3),
                                   f"{int(scores[i]*100)}%", font_size=Pt(10), color=WHITE,
                                   alignment=PP_ALIGN.RIGHT)

    # Right: radar chart mockup
    radar_x = Inches(7.8)
    radar_y = Inches(1.5)
    radar_size = Inches(3.5)
    center_x = radar_x + radar_size / 2
    center_y = radar_y + radar_size / 2
    radius = Inches(1.4)

    # Draw hexagon background
    hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, radar_x, radar_y, radar_size, radar_size)
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = SURFACE_CONTAINER
    hex_shape.line.color.rgb = OUTLINE_VARIANT
    hex_shape.line.width = Pt(1)

    # Inner hexagon (semi-transparent not possible in python-pptx, use lighter fill)
    inner_hex = slide.shapes.add_shape(MSO_SHAPE.HEXAGON,
                                        radar_x + Inches(0.7), radar_y + Inches(0.7),
                                        radar_size - Inches(1.4), radar_size - Inches(1.4))
    inner_hex.fill.solid()
    inner_hex.fill.fore_color.rgb = RGBColor(0xD8, 0xE2, 0xFF)  # primary-fixed
    inner_hex.line.color.rgb = PRIMARY
    inner_hex.line.width = Pt(1.5)

    # Data points on hexagon vertices (6 vertices at 60° intervals)
    import math
    dim_labels = ["睡眠", "起床", "吸烟", "饮酒", "噪音", "清洁"]
    data_scores = [0.85, 0.78, 0.92, 0.88, 0.75, 0.80]
    cx_pt = radar_x + radar_size / 2
    cy_pt = radar_y + radar_size / 2

    for i in range(6):
        angle = -math.pi / 2 + i * math.pi / 3  # start from top, go clockwise
        r = radius * data_scores[i]
        dx = Emu(r) * math.cos(angle)
        dy = Emu(r) * math.sin(angle)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      int(cx_pt) + dx - Inches(0.1),
                                      int(cy_pt) + dy - Inches(0.1),
                                      Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = PRIMARY
        dot.line.fill.background()
        # Label
        label_r = radius + Inches(0.25)
        lx = int(cx_pt) + Emu(label_r) * math.cos(angle)
        ly = int(cy_pt) + Emu(label_r) * math.sin(angle)
        add_text_box(slide, lx - Inches(0.35), ly - Inches(0.15), Inches(0.7), Inches(0.3),
                     dim_labels[i], font_size=Pt(9), color=ON_SURFACE_VARIANT, alignment=PP_ALIGN.CENTER)

    # Caption
    add_text_box(slide, radar_x, radar_y + radar_size + Inches(0.4), radar_size, Inches(0.3),
                 "▲ DeepSeek 六维匹配雷达图（基于习惯问卷相似度）",
                 font_size=Pt(9), color=OUTLINE, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 7)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 8: AI CHATBOT
# ══════════════════════════════════════════════════════════════════════
def slide_ai_chat():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "AI 客服助手", "AI Customer Service Chatbot")

    # Left: system info
    left_x = Inches(0.7)
    add_card(slide, left_x, Inches(1.2), Inches(5.5), Inches(1.8),
             "系统角色定位",
             [("• 充当宿舍管理智能助手，回答学生关于宿舍设施、管理规定、分配流程等问题", False, Pt(12)),
              ("• 基于学生个人上下文（姓名、专业、年级、宿舍、习惯问卷）提供个性化建议", False, Pt(12)),
              ("• 支持宿舍调换政策咨询、报修流程指导、生活习惯适配建议等", False, Pt(12))],
             accent_color=SECONDARY, title_size=Pt(15))

    add_card(slide, left_x, Inches(3.25), Inches(5.5), Inches(2.5),
             "技术实现",
             [("ChatController — AJAX 异步请求，JSON 格式响应", False, Pt(11)),
              ("AIServiceImpl — 调用 DeepSeek Chat API (deepseek-chat)", False, Pt(11)),
              ("HttpSession — 会话级对话历史存储（最多 40 条消息）", False, Pt(11)),
              ("参数配置 — max_tokens=2048, temperature=0.7", False, Pt(11)),
              ("前端交互 — 消息气泡、输入中状态指示、一键清空历史", False, Pt(11))],
             accent_color=SECONDARY, title_size=Pt(15))

    # Right: mock chat bubbles
    chat_x = Inches(6.7)
    chat_y = Inches(1.2)

    # Chat container background
    add_rounded_rect(slide, chat_x, chat_y, Inches(5.8), Inches(4.55),
                     fill_color=SURFACE_LOWEST, border_color=OUTLINE_VARIANT)

    # Chat header
    add_rect(slide, chat_x, chat_y, Inches(5.8), Inches(0.45), fill_color=SECONDARY)
    add_text_box(slide, chat_x + Inches(0.15), chat_y + Inches(0.05), Inches(5.5), Inches(0.35),
                 "AI 客服对话", font_size=Pt(11), bold=True, color=WHITE)

    # AI message bubble (left)
    ai_bubble_x = chat_x + Inches(0.2)
    ai_bubble_y = chat_y + Inches(0.7)
    ai_bubble = add_rounded_rect(slide, ai_bubble_x, ai_bubble_y, Inches(3.8), Inches(0.9),
                                 fill_color=SURFACE_CONTAINER)
    add_text_box(slide, ai_bubble_x + Inches(0.15), ai_bubble_y + Inches(0.08), Inches(3.5), Inches(0.2),
                 "AI 助手", font_size=Pt(9), bold=True, color=SECONDARY)
    add_text_box(slide, ai_bubble_x + Inches(0.15), ai_bubble_y + Inches(0.30), Inches(3.5), Inches(0.5),
                 "你好！我是宿舍管理助手，有什么可以帮你？",
                 font_size=Pt(10), color=ON_SURFACE)

    # User message bubble (right)
    user_bubble_x = chat_x + Inches(2.0)
    user_bubble_y = ai_bubble_y + Inches(1.15)
    user_bubble = add_rounded_rect(slide, user_bubble_x, user_bubble_y, Inches(3.4), Inches(0.7),
                                   fill_color=RGBColor(0xD8, 0xE2, 0xFF))
    add_text_box(slide, user_bubble_x + Inches(0.15), user_bubble_y + Inches(0.08), Inches(3.1), Inches(0.2),
                 "学生", font_size=Pt(9), bold=True, color=PRIMARY)
    add_text_box(slide, user_bubble_x + Inches(0.15), user_bubble_y + Inches(0.30), Inches(3.1), Inches(0.4),
                 "请问怎么申请调换宿舍？",
                 font_size=Pt(10), color=ON_SURFACE)

    # AI reply bubble
    ai_reply_y = user_bubble_y + Inches(0.9)
    ai_reply = add_rounded_rect(slide, ai_bubble_x, ai_reply_y, Inches(4.2), Inches(1.1),
                                fill_color=SURFACE_CONTAINER)
    add_text_box(slide, ai_bubble_x + Inches(0.15), ai_reply_y + Inches(0.08), Inches(3.9), Inches(0.2),
                 "AI 助手", font_size=Pt(9), bold=True, color=SECONDARY)
    add_text_box(slide, ai_bubble_x + Inches(0.15), ai_reply_y + Inches(0.30), Inches(3.9), Inches(0.7),
                 "你可以在「宿舍调换」页面提交申请，填写调换原因并上传附件，管理员审核后会通知你。",
                 font_size=Pt(10), color=ON_SURFACE)

    # Input bar
    input_y = chat_y + Inches(3.85)
    add_rounded_rect(slide, chat_x + Inches(0.15), input_y, Inches(4.8), Inches(0.38),
                     fill_color=SURFACE, border_color=OUTLINE_VARIANT)
    add_text_box(slide, chat_x + Inches(0.3), input_y + Inches(0.04), Inches(4.5), Inches(0.3),
                 "输入消息...", font_size=Pt(10), color=OUTLINE)

    add_page_number(slide, 8)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 9: DATABASE DESIGN
# ══════════════════════════════════════════════════════════════════════
def slide_database():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "数据库设计", "Database Schema")

    # ER diagram - 7 tables arranged in three rows
    tables = [
        ("user", "用户表", ["id", "username", "password", "role", "status"]),
        ("student", "学生表", ["id", "user_id(FK)", "name", "gender", "major", "grade", "student_id"]),
        ("dormitory", "宿舍表", ["id", "dormitory_number", "dormitory_type", "capacity", "current_occupancy", "building"]),
        ("habit_questionnaire", "习惯问卷表", ["id", "student_id(FK)", "sleep_time", "wake_up_time", "smoking", "drinking", "noise_tolerance", "cleanliness"]),
        ("dormitory_allocation", "分配表", ["id", "student_id(FK)", "dormitory_id(FK)", "allocate_time", "status"]),
        ("dormitory_transfer", "调换表", ["id", "student_id(FK)", "current_dormitory_id", "target_dormitory_id", "reason", "attachment", "status"]),
        ("repair_request", "报修表", ["id", "student_id(FK)", "dormitory_id(FK)", "description", "image", "status"]),
    ]

    # Layout: 1st row (1), 2nd row (3), 3rd row (3)
    row_configs = [
        (1, [0]),       # user center top
        (3, [1, 2, 3]), # student, dormitory, questionnaire
        (3, [4, 5, 6]), # allocation, transfer, repair
    ]

    table_h = Inches(1.35)
    table_w = Inches(3.6)
    y_start_db = Inches(1.05)

    positions = []  # (table_name, cx, cy, width, height)

    for row_idx, (count, indices) in enumerate(row_configs):
        total_w = count * table_w + (count - 1) * Inches(0.25)
        start_x = (SLIDE_W - total_w) / 2
        y = y_start_db + row_idx * (table_h + Inches(0.35))
        for j, table_idx in enumerate(indices):
            x = start_x + j * (table_w + Inches(0.25))
            eng_name, cn_name, fields = tables[table_idx]
            accent = PRIMARY if table_idx == 0 else (SECONDARY if table_idx in [1, 3] else ON_SURFACE)

            add_card(slide, x, y, table_w, table_h, f"{cn_name} ({eng_name})",
                     [(f, False, Pt(9)) for f in fields],
                     accent_color=accent, title_size=Pt(11))
            positions.append((eng_name, x + table_w / 2, y, table_w, table_h))

    # Connectors between related tables
    # user -> student
    user_pos = positions[0]
    student_pos = positions[1]
    add_connector(slide, user_pos[1], user_pos[2] + table_h,
                  student_pos[1], student_pos[2], OUTLINE_VARIANT)

    # Bottom note
    add_text_box(slide, Inches(0.7), SLIDE_H - Inches(0.65), Inches(11.9), Inches(0.3),
                 "数据库：MySQL 8.0  |  连接池：Druid 1.2.22  |  ORM：MyBatis 3.5.16  |  共 7 张业务表，完整外键约束与索引",
                 font_size=Pt(10), color=ON_SURFACE_VARIANT, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 9)


# ══════════════════════════════════════════════════════════════════════
# SLIDE 10: SUMMARY
# ══════════════════════════════════════════════════════════════════════
def slide_summary():
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, SURFACE)
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), fill_color=PRIMARY)

    add_accent_header(slide, Inches(0.7), Inches(0.35), Inches(11), "项目亮点与总结", "Highlights & Summary")

    # 2x2 highlights grid
    highlights = [
        ("AI 赋能管理", "",
         "集成 DeepSeek 大模型，实现\n智能宿舍推荐与 24×7 在线客服\n大幅提升管理效率与用户体验",
         PRIMARY),
        ("数据驱动决策", "",
         "ECharts 雷达图六维可视化\n基于生活习惯问卷的科学匹配\n为宿舍分配提供量化依据",
         SECONDARY),
        ("全流程闭环", "",
         "覆盖入住分配 → 调换审批 → 报修处理\n全业务流程，状态流转清晰\n支持文件附件与图片上传",
         PRIMARY),
        ("现代技术栈", "",
         "Java 21 + Spring 6 + MyBatis\nMD3 设计语言 + Tailwind CSS\nMaven 构建 + Tomcat 容器化部署",
         SECONDARY),
    ]

    card_w6 = Inches(5.7)
    card_h6 = Inches(2.35)
    y6 = Inches(1.15)
    gap6 = Inches(0.25)

    for i, (title, icon, desc, accent) in enumerate(highlights):
        row, col = i // 2, i % 2
        x = Inches(0.7) + col * (card_w6 + gap6)
        y = y6 + row * (card_h6 + gap6)

        add_card(slide, x, y, card_w6, card_h6, title, [], accent_color=accent, title_size=Pt(17))
        add_text_box(slide, x + card_w6 - Inches(0.85), y + Inches(0.15), Inches(0.7), Inches(0.7),
                     icon, font_size=Pt(36), color=OUTLINE_VARIANT, font_name=FONT_EN,
                     anchor=MSO_ANCHOR.MIDDLE)
        add_multiline_box(slide, x + Inches(0.2), y + Inches(0.6), card_w6 - Inches(0.35),
                          card_h6 - Inches(0.7),
                          [(line, False, Pt(13)) for line in desc.split("\n")],
                          font_size=Pt(13), color=ON_SURFACE_VARIANT)

    # Tech stack tag cloud
    tags = ["Java 21", "Spring 6", "Spring MVC", "MyBatis", "MySQL", "Druid", "JSP", "Tailwind CSS",
            "ECharts", "DeepSeek", "Apache HttpClient", "FastJSON", "Maven", "Tomcat", "Log4j"]
    tag_y = Inches(6.15)
    tag_start_x = Inches(1.5)
    cur_x = tag_start_x
    for tag in tags:
        tag_w = Inches(0.3 + len(tag) * 0.13)
        if cur_x + tag_w > SLIDE_W - Inches(0.8):
            cur_x = tag_start_x
            tag_y += Inches(0.35)
        add_badge(slide, cur_x, tag_y, tag, color=OUTLINE)
        cur_x += tag_w + Inches(0.12)

    # Thank you
    add_text_box(slide, Inches(0), SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.35),
                 "感谢观看", font_size=Pt(14), color=ON_SURFACE_VARIANT, alignment=PP_ALIGN.CENTER)

    add_page_number(slide, 10)


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════
def main():
    slide_cover()
    slide_overview()
    slide_architecture()
    slide_modules()
    slide_admin()
    slide_student()
    slide_ai_recommend()
    slide_ai_chat()
    slide_database()
    slide_summary()

    output_path = os.path.join(os.path.dirname(__file__), "宿舍管理系统-项目展示.pptx")
    prs.save(output_path)
    print(f"PPT generated: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
